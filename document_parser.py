"""Extract plain text from various document formats."""

from pathlib import Path


def extract_text(file_path: str | Path) -> dict:
    """Extract text from a file based on its extension.

    Returns dict with: text, title, word_count, doc_type, source_url
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    extractors = {
        ".md": _extract_markdown,
        ".txt": _extract_markdown,
        ".docx": _extract_docx,
        ".pdf": _extract_pdf,
        ".pptx": _extract_pptx,
    }

    extractor = extractors.get(ext, _extract_markdown)
    text = extractor(path)
    title = _infer_title(text, path)

    return {
        "text": text,
        "title": title,
        "word_count": len(text.split()),
        "doc_type": ext.lstrip(".") or "text",
        "source_url": "",
    }


def _extract_markdown(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[Error extracting .docx: {e}]"


def _extract_pdf(path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except Exception as e:
        return f"[Error extracting .pdf: {e}]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[Error extracting .pptx: {e}]"


def _infer_title(text: str, path: Path) -> str:
    """Try to get title from first heading, fall back to filename."""
    for line in text.split("\n"):
        line = line.strip()
        if line.startswith("# "):
            return line[2:].strip()
        if line and not line.startswith("#"):
            break
    return path.stem.replace("-", " ").replace("_", " ").title()
